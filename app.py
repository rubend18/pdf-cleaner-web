import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
import io
import os
import numpy as np

st.set_page_config(page_title="PDF a PPTX", layout="centered")

st.title("📄 PDF a PPTX")
st.write("Sube tu PDF para tapar el logo de NotebookLM y convertirlo a PowerPoint.")

uploaded_file = st.file_uploader("", type=["pdf"], label_visibility="collapsed")

def procesar_pdf(file_stream):
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    prs = Presentation()

    # Configuración del recuadro
    ancho_box = 97
    alto_box = 12
    margen_inf = 9
    margen_der = 10

    total_pages = len(doc)
    progress_bar = st.progress(0)

    for i, page in enumerate(doc):
        w_page = page.rect.width
        h_page = page.rect.height

        rect = fitz.Rect(
            w_page - margen_der - ancho_box,
            h_page - margen_inf - alto_box,
            w_page - margen_der,
            h_page - margen_inf
        )

        # Detección de color promedio
        try:
            pix_crop = page.get_pixmap(clip=rect)
            img_data = np.frombuffer(pix_crop.samples, dtype=np.uint8)
            if len(img_data) > 0:
                img_data = img_data.reshape(pix_crop.h, pix_crop.w, pix_crop.n)
                avg_color = img_data.mean(axis=(0, 1))
                r, g, b = avg_color[:3] / 255.0
            else:
                r, g, b = 1, 1, 1
        except:
            r, g, b = 1, 1, 1

        # Dibujar parche
        page.draw_rect(rect, color=(r, g, b), fill=(r, g, b))

        # Renderizar y crear slide
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_bytes = pix.tobytes("png")
        image_stream = io.BytesIO(img_bytes)

        prs.slide_width = Pt(w_page)
        prs.slide_height = Pt(h_page)
        blank_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

        progress_bar.progress((i + 1) / total_pages)

    pptx_out = io.BytesIO()
    prs.save(pptx_out)
    pptx_out.seek(0)
    return pptx_out

if uploaded_file is not None:
    # Procesar automáticamente al cargar o con un botón simple
    if st.button("Convertir"):
        with st.spinner("Procesando..."):
            try:
                resultado = procesar_pdf(uploaded_file)
                st.download_button(
                    label="Descargar PPTX",
                    data=resultado,
                    file_name=f"{os.path.splitext(uploaded_file.name)[0]}_ppt.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception:
                st.error("Error al procesar el archivo.")